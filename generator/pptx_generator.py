"""
PowerPoint-Generator
Ungeverändert aus WBI-Docu-Assist übernommen.
"""

import io
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _safe_filename(title: str) -> str:
    return re.sub(r'[^\w\s\-äöüÄÖÜß]', '_', title).strip() + ".pptx"


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color='333333',
                  align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


def _get_slide_layout(prs, name_contains, fallback_idx=1):
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout
    return prs.slide_layouts[min(fallback_idx, len(prs.slide_layouts) - 1)]


def generate_pptx(data: dict, template_path: str):
    title = data.get('titleSubject', '') + ' - ' + data.get('titleTopic', '')
    title = title.strip(' -')
    chapters = data.get('chapters', [])
    valid_chapters = [c for c in chapters if c.get('name', '').strip()]
    today = datetime.now().strftime('%d.%m.%Y')

    prs = Presentation(template_path)
    W = prs.slide_width
    H = prs.slide_height
    MARGIN = Inches(0.6)
    CONTENT_W = W - 2 * MARGIN
    BLUE = '1F3864'
    WHITE = 'FFFFFF'
    LBLUE = 'BDD7EE'

    def add_slide(layout_hint='Inhalt', fallback=1):
        layout = _get_slide_layout(prs, layout_hint, fallback)
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            try:
                ph.text = ''
            except Exception:
                pass
        return slide

    # Titelfolie
    title_slide = add_slide('Titel', 0)
    for shape in list(title_slide.shapes):
        title_slide.shapes._spTree.remove(shape._element)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(BLUE)
    _add_text_box(title_slide, title, MARGIN, H // 3, CONTENT_W, Inches(1.5),
                  font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(title_slide, today, MARGIN, H // 3 + Inches(1.8), CONTENT_W, Inches(0.5),
                  font_size=14, color=LBLUE, align=PP_ALIGN.CENTER)

    # Agenda-Folie
    if valid_chapters:
        agenda_slide = add_slide('Inhalt', 1)
        for shape in list(agenda_slide.shapes):
            agenda_slide.shapes._spTree.remove(shape._element)
        hdr = agenda_slide.shapes.add_shape(1, 0, 0, W, Inches(1.0))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _rgb(BLUE)
        hdr.line.fill.background()
        _add_text_box(agenda_slide, 'Agenda', MARGIN, Inches(0.12), CONTENT_W, Inches(0.75),
                      font_size=26, bold=True, color=WHITE)
        agenda_text = '\n'.join(f'{i+1}.  {c["name"]}' for i, c in enumerate(valid_chapters))
        _add_text_box(agenda_slide, agenda_text, MARGIN, Inches(1.1), CONTENT_W, H - Inches(1.5),
                      font_size=16, color='333333')

    # Inhaltsfolien
    for ch in valid_chapters:
        content_slide = add_slide('Inhalt', 1)
        for shape in list(content_slide.shapes):
            content_slide.shapes._spTree.remove(shape._element)
        hdr = content_slide.shapes.add_shape(1, 0, 0, W, Inches(1.0))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _rgb(BLUE)
        hdr.line.fill.background()
        _add_text_box(content_slide, ch['name'], MARGIN, Inches(0.12), CONTENT_W, Inches(0.75),
                      font_size=22, bold=True, color=WHITE)
        subs = [s for s in ch.get('subs', []) if s.strip()]
        content = '\n'.join(f'•  {s}' for s in subs) if subs else 'Inhalt hier einfügen.'
        _add_text_box(content_slide, content, MARGIN, Inches(1.15), CONTENT_W, H - Inches(1.6),
                      font_size=16, color='333333', italic=not bool(subs))

    # Abschlussfolie
    end_slide = add_slide('Titel', 0)
    for shape in list(end_slide.shapes):
        end_slide.shapes._spTree.remove(shape._element)
    bg2 = end_slide.background
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = _rgb(BLUE)
    _add_text_box(end_slide, 'Vielen Dank', MARGIN, H // 3, CONTENT_W, Inches(1.5),
                  font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf, _safe_filename(title)
